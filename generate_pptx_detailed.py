import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Please install it using pip.")
    sys.exit(1)

def generate_detailed_pptx():
    print("Initializing Overhauled PowerPoint Presentation generation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen layout
    prs.slide_height = Inches(7.5)
    
    # Theme color definitions
    c_bg = RGBColor(10, 15, 29)       # Deep Dark Navy #0A0F1D
    c_title = RGBColor(248, 250, 252) # White #F8FAFC
    c_body = RGBColor(148, 163, 184)  # Muted Slate #94A3B8
    c_blue = RGBColor(96, 165, 250)   # Indigo/Blue #60A5FA

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    slides_content = [
        # Slide 1: Title
        {
            "is_title": True,
            "title": "SECURE AICTE MEETING & GOVERNANCE PLATFORM",
            "subtitle": "Complete Overhaul & Integrated Multi-Module Security Suite\n\nWalkthrough Deck: Light/Dark Mode, MFA, Jitsi WebRTC, S3 Files, Blockchain, & AI MoMs\nSmart India Hackathon Presentation Deck"
        },
        # Slide 2: Global UI/UX Design System
        {
            "title": "Topic: Global UI/UX Design & Theme System",
            "sections": [
                {
                    "heading": "Light & Dark Mode Integration",
                    "text": "• Responsive CSS Variables: Configured in index.css with mapped root custom properties.\n"
                         "• LocalStorage Persistence: User theme selections remain consistent on page reloads.\n"
                         "• System Preference Checks: Respects browser media queries on first visit.\n"
                         "• Theme Toggler: Polished rotating Sun/Moon button animation in the primary header."
                },
                {
                    "heading": "Premium Aesthetics",
                    "text": "• Colors: Deep dark navy background (#0A0F1D) in dark mode, white/slate slate (#F8FAFC) in light mode.\n"
                         "• Glassmorphism: Cards feature glass-filtered backdrops with subtle border details.\n"
                         "• Micro-Animations: Modal fade-ins, card scaling, and warning pulse glow states."
                }
            ]
        },
        # Slide 3: Authentication & MFA (Module 1)
        {
            "title": "Module 1: Authentication & MFA Gate",
            "sections": [
                {
                    "heading": "Multi-Factor Authentication Flow",
                    "text": "• Username/Password Validation: Primary login credentials check (admin_aicte / password123).\n"
                         "• Passcode Toggle: Toggles visibility of the secure password access key.\n"
                         "• OTP Dialog: Displays 6-digit dynamic authenticator (TOTP) passcode prompt."
                },
                {
                    "heading": "Session Security Control",
                    "text": "• Access Token Sign-out: Dropdown in header terminates sessions and clears memory states.\n"
                         "• RBAC Access: Displays officer credential badges based on role-based access control."
                }
            ]
        },
        # Slide 4: Meetings & Jitsi WebRTC (Module 2)
        {
            "title": "Module 2: Council Meetings & Jitsi Meet",
            "sections": [
                {
                    "heading": "WebRTC Video Conferencing",
                    "text": "• Jitsi External API: Dynamically fetches external_api.js script to mount active conference iframes.\n"
                         "• Audio/Video Hardware Toggle: Custom controls in room control bar interface."
                },
                {
                    "heading": "Governance Scheduling",
                    "text": "• Schedule Modal: Configures meeting details with a security sensitivity level picker (Low to Top Secret).\n"
                         "• Security Badge: Displays active E2EE encryption status in meeting rooms."
                }
            ]
        },
        # Slide 5: Meeting Recordings Archive (Module 2/3)
        {
            "title": "Module 2/3: Secure Video Archive Portal",
            "sections": [
                {
                    "heading": "Interactive Media Player",
                    "text": "• Custom Controls: Features custom Play, Pause, Time Seek bar, and Volume controls.\n"
                         "• Object URL Ingestion: Enables dragging and dropping MP4/WebM recordings to play them locally."
                },
                {
                    "heading": "Audit Trails & Watermarks",
                    "text": "• Security Overlay: Watermarks confidentiality text with user credentials and timestamp.\n"
                         "• Integrity Seal: Displays SHA-256 hashes and access clearance role listings."
                }
            ]
        },
        # Slide 6: Secure S3 Document Vault (Module 3)
        {
            "title": "Module 3: Secure S3 File Management",
            "sections": [
                {
                    "heading": "Secure Upload System",
                    "text": "• Drag & Drop Zone: Drag PDF reports or Excel sheets to trigger upload sequence animations.\n"
                         "• File Validation: Check size restrictions (max 20MB) and formats before hashing."
                },
                {
                    "heading": "Storage Audits & Metadata",
                    "text": "• S3 Indicators: Displays live MinIO/S3 cluster status.\n"
                         "• File Clearance Inspector: View AES-256 encryption tags, SHA-256 hashes, and version counts."
                }
            ]
        },
        # Slide 7: Blockchain Verification (Module 4)
        {
            "title": "Module 4: Blockchain Ledger Verification",
            "sections": [
                {
                    "heading": "Integrity Check Terminal",
                    "text": "• SHA-256 Verifier: Paste file hashes to query the blockchain registry.\n"
                         "• Integrity Statuses: Green verified organizational badge or Red compromised warning panels."
                },
                {
                    "heading": "Fabric Block Explorer",
                    "text": "• Visual Timelines: Displays newly committed ledger blocks.\n"
                         "• Metadata: Tracks transaction IDs, channels, block heights, and anchoring stamps."
                }
            ]
        },
        # Slide 8: AI Meeting Intelligence (Module 5)
        {
            "title": "Module 5: AI Meeting Intelligence Engine",
            "sections": [
                {
                    "heading": "AI Transcripts & Summarization",
                    "text": "• Speaker Tagged Feed: Automatically splits spoken dialogue with exact timestamps.\n"
                         "• MoM Minutes: Compiles goals summaries, approvals, and decisions list."
                },
                {
                    "heading": "Human Sign-off Workflow",
                    "text": "• Review Portal: Action items list tracks assignees, deadlines, and progress states.\n"
                         "• Sign-off Steps: Transitions records through three states (AI-Generated ➔ Human Approved ➔ Final Commit)."
                }
            ]
        },
        # Slide 9: Audit Logs & Attendance Tracker (Module 6)
        {
            "title": "Module 6: Compliance Auditing & Attendance",
            "sections": [
                {
                    "heading": "Attendance Duration Logs",
                    "text": "• Jitsi Webhooks: Captures user connect/disconnect timestamps.\n"
                         "• Duration Formula: Computes total minutes, flagging Late and Left Early parameters."
                },
                {
                    "heading": "System Activity Auditing",
                    "text": "• M6 Ledger: Logs user actions, modules, origins, IP, and severities.\n"
                         "• SOC Alarms: Quarantines threat origins and prompts mitigation actions."
                }
            ]
        },
        # Slide 10: Stack & Widescreen Deliverables
        {
            "title": "Widescreen Prototype Architecture",
            "sections": [
                {
                    "heading": "Front-End UI Deliverables",
                    "text": "React + Vite widescreen framework. Widescreen widescreen layout (16:9 widescreen layout), Tailwind CSS v4 design system, and Lucide vector icons."
                },
                {
                    "heading": "Mock Server API & Databases",
                    "text": "Express API server with mock databases. Client-side state fallback coordinates operations offline. Linked GitHub codebase and live Vercel cloud deploy."
                }
            ]
        }
    ]

    for idx, slide_data in enumerate(slides_content):
        slide = prs.slides.add_slide(blank_layout)
        
        # Slide Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg
        
        if slide_data.get("is_title"):
            # Title slide layout
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            p1.text = slide_data["title"]
            p1.font.size = Pt(40)
            p1.font.bold = True
            p1.font.color.rgb = c_title
            p1.font.name = 'Arial'
            
            p2 = tf.add_paragraph()
            p2.text = "\n" + slide_data["subtitle"]
            p2.font.size = Pt(16)
            p2.font.color.rgb = c_blue
            p2.font.name = 'Arial'
            
        else:
            # Header
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
            tf_h = header_box.text_frame
            tf_h.word_wrap = True
            p_h = tf_h.paragraphs[0]
            p_h.text = slide_data["title"]
            p_h.font.size = Pt(32)
            p_h.font.bold = True
            p_h.font.color.rgb = c_blue
            p_h.font.name = 'Arial'
            
            # Sub-sections (Detailed left & right boxes)
            sec1 = slide_data["sections"][0]
            sec2 = slide_data["sections"][1]
            
            # Left block
            left_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.8))
            tf_l = left_box.text_frame
            tf_l.word_wrap = True
            
            p_l_h = tf_l.paragraphs[0]
            p_l_h.text = sec1["heading"]
            p_l_h.font.size = Pt(20)
            p_l_h.font.bold = True
            p_l_h.font.color.rgb = c_title
            p_l_h.font.name = 'Arial'
            p_l_h.space_after = Pt(10)
            
            p_l_b = tf_l.add_paragraph()
            p_l_b.text = sec1["text"]
            p_l_b.font.size = Pt(14)
            p_l_b.font.color.rgb = c_body
            p_l_b.font.name = 'Arial'
            p_l_b.line_spacing = 1.2
            
            # Right block
            right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.4), Inches(4.8))
            tf_r = right_box.text_frame
            tf_r.word_wrap = True
            
            p_r_h = tf_r.paragraphs[0]
            p_r_h.text = sec2["heading"]
            p_r_h.font.size = Pt(20)
            p_r_h.font.bold = True
            p_r_h.font.color.rgb = c_title
            p_r_h.font.name = 'Arial'
            p_r_h.space_after = Pt(10)
            
            p_r_b = tf_r.add_paragraph()
            p_r_b.text = sec2["text"]
            p_r_b.font.size = Pt(14)
            p_r_b.font.color.rgb = c_body
            p_r_b.font.name = 'Arial'
            p_r_b.line_spacing = 1.2

            # Footer
            footer_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.8), Inches(1.5), Inches(0.5))
            tf_f = footer_box.text_frame
            p_f = tf_f.paragraphs[0]
            p_f.text = f"Slide {idx + 1} / 10"
            p_f.font.size = Pt(10)
            p_f.font.color.rgb = c_body
            p_f.font.name = 'Arial'

    prs.save("AICTE_M6_Presentation_Detailed.pptx")
    print("Overhauled PowerPoint updated successfully as 'AICTE_M6_Presentation_Detailed.pptx'")

if __name__ == "__main__":
    generate_detailed_pptx()
