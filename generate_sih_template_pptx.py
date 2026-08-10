import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Please install it using pip.")
    sys.exit(1)

def create_sih_presentation():
    print("Initializing SIH 2026 Idea Submission Template PowerPoint generation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen layout
    prs.slide_height = Inches(7.5)
    
    # Theme color definitions
    c_white = RGBColor(255, 255, 255)
    c_black = RGBColor(0, 0, 0)
    c_blue_banner = RGBColor(25, 118, 210) # SIH bottom banner blue
    c_heading_blue = RGBColor(19, 78, 175) # Deep heading blue
    c_muted = RGBColor(100, 116, 139)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    slides_data = [
        # Slide 1: TITLE PAGE
        {
            "slide_num": 1,
            "type": "title",
            "title": "TITLE PAGE",
            "bullets": [
                "Problem Statement ID – SIH1608 (AICTE Governance)",
                "Problem Statement Title – Secure AICTE Online Meeting & Governance Platform",
                "Theme – Smart Automation / Security",
                "PS Category – Software",
                "Team ID – [Enter Your Team ID]",
                "Team Name (Registered on portal) – [Enter Registered Team Name]"
            ]
        },
        # Slide 2: PROPOSED SOLUTION
        {
            "slide_num": 2,
            "type": "solution",
            "title": "SECURE AICTE MEETING & GOVERNANCE PLATFORM",
            "subtitle": "❖ Proposed Solution (Describe your Idea/Solution/Prototype)",
            "bullets": [
                "Centralized Auditing & Governance (M6): Ingests audits across M1 (Auth), M2 (Meetings), M3 (Files), M4 (Ledger), and M5 (AI Transcripts), acting as a centralized logs database.",
                "E2EE Video Conferencing & S3 Storage: Features Jitsi WebRTC video call frames in portal alongside S3 upload modules equipped with automated AES-256 encryption.",
                "Local Video Object Streaming: Solves database size constraints by enabling users to upload and stream meeting recordings locally in the browser via Object URLs.",
                "Blockchain Anchor Locks: Encodes PDF reports to SHA-256 hashes, committing transaction tokens to private blockchain blocks to ensure immutable records."
            ]
        },
        # Slide 3: TECHNICAL APPROACH
        {
            "slide_num": 3,
            "type": "standard",
            "title": "TECHNICAL APPROACH",
            "bullets": [
                "Technologies to be used:",
                "  • Frontend Client: React 19, Vite, Tailwind CSS v4, Recharts Analytics, Lucide vector icons.",
                "  • Backend Services: Node.js + Express, Jitsi WebRTC External API script integration.",
                "  • Cryptography & Databases: PostgreSQL schemas, SHA-256 hashing, and AES-256-CBC encryption.",
                "Methodology and process for implementation:",
                "  • Authentication: Initial secure credentials check triggers 6-digit MFA OTP validation.",
                "  • Meeting & AI: Post-call transcripts feed speaker-split segments to AI summary engines.",
                "  • Sealing: Officer review approves AI draft MoMs, committing report hash tokens to blockchain ledger."
            ]
        },
        # Slide 4: FEASIBILITY AND VIABILITY
        {
            "slide_num": 4,
            "type": "standard",
            "title": "FEASIBILITY AND VIABILITY",
            "bullets": [
                "Analysis of the feasibility of the idea:",
                "  • Feasible deployment: Connects open-source WebRTC (Jitsi) directly with lightweight mock database states.",
                "  • Administrative integration: Integrates easily with existing technical college governance frameworks.",
                "Potential challenges and risks:",
                "  • Database bloating due to video files, and internet outages during live video streaming.",
                "Strategies for overcoming these challenges:",
                "  • Client-side fallback: Enables local memory fallback states if PostgreSQL node is offline.",
                "  • Object URL playback: Videos are uploaded and streamed locally, bypassing backend storage costs."
            ]
        },
        # Slide 5: IMPACT AND BENEFITS
        {
            "slide_num": 5,
            "type": "standard",
            "title": "IMPACT AND BENEFITS",
            "bullets": [
                "Potential impact on the target audience:",
                "  • Prevents administrative tampering: Assures educational institutions of the integrity of directive decisions.",
                "  • Streamlines audits: Compliance officers check logs and resolve threat levels through a single dashboard.",
                "Benefits of the solution (social, economic, environmental, etc.):",
                "  • Minimizes paper printouts and courier requirements by archiving signed PDF reports securely.",
                "  • Avoids expensive subscriptions by deploying free, open-source secure WebRTC (Jitsi) platforms.",
                "  • Accelerates national hackathon evaluations by deploying automated AI summarizers."
            ]
        },
        # Slide 6: RESEARCH AND REFERENCES
        {
            "slide_num": 6,
            "type": "standard",
            "title": "RESEARCH AND REFERENCES",
            "bullets": [
                "Details / Links of the reference and research work:",
                "  • Jitsi Meet External API Integrations: https://jitsi.github.io/handbook/docs/dev-guide/dev-guide-iframe/",
                "  • Hyperledger Fabric Channel Architecture & Chaincodes: https://hyperledger-fabric.readthedocs.io/",
                "  • W3C WebRTC (Web Real-Time Communication) standards for E2EE audio/video streaming.",
                "  • PostgreSQL Full-Text Search (FTS) Lexeme parser & pgvector Indexing methods.",
                "  • AES-256-GCM symmetric encryption and SHA-256 cryptographic hash standard specifications."
            ]
        }
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background to solid white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_white

        # 1. Top Left Circle representing "Your Team Name"
        team_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(0.5), Inches(0.3), Inches(1.2), Inches(1.2)
        )
        team_circle.fill.solid()
        team_circle.fill.fore_color.rgb = c_white
        team_circle.line.color.rgb = RGBColor(120, 80, 160) # Purple circle border
        team_circle.line.width = Pt(1.5)
        
        tf_tc = team_circle.text_frame
        tf_tc.word_wrap = True
        p_tc = tf_tc.paragraphs[0]
        p_tc.text = "Your\nTeam\nName"
        p_tc.font.size = Pt(11)
        p_tc.font.color.rgb = c_black
        p_tc.font.name = 'Arial'
        p_tc.alignment = PP_ALIGN.CENTER

        # 2. Top Right Logo Representation "SMART INDIA HACKATHON 2026"
        logo_box = slide.shapes.add_textbox(Inches(10.2), Inches(0.2), Inches(2.8), Inches(1.3))
        tf_lb = logo_box.text_frame
        tf_lb.word_wrap = True
        p_lb1 = tf_lb.paragraphs[0]
        p_lb1.text = "SMART INDIA"
        p_lb1.font.size = Pt(13)
        p_lb1.font.bold = True
        p_lb1.font.color.rgb = c_heading_blue
        p_lb1.font.name = 'Arial'
        p_lb1.alignment = PP_ALIGN.RIGHT
        
        p_lb2 = tf_lb.add_paragraph()
        p_lb2.text = "HACKATHON 2026"
        p_lb2.font.size = Pt(12)
        p_lb2.font.bold = True
        p_lb2.font.color.rgb = c_heading_blue
        p_lb2.font.name = 'Arial'
        p_lb2.alignment = PP_ALIGN.RIGHT

        # 3. Bottom Blue Footer Banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.9), Inches(13.333), Inches(0.6)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = c_blue_banner
        banner.line.fill.background() # No border line

        # Footer Text inside Banner
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.5))
        tf_fb = footer_box.text_frame
        p_fb = tf_fb.paragraphs[0]
        p_fb.text = f"@SIH Idea submission- Template                                                                                                                                                 {slide_data['slide_num']}"
        p_fb.font.size = Pt(11)
        p_fb.font.bold = True
        p_fb.font.color.rgb = c_white
        p_fb.font.name = 'Arial'

        # 4. Content Formatting based on type
        if slide_data["type"] == "title":
            # Header title (centered top)
            header_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.3), Inches(8.0), Inches(0.8))
            tf_hb = header_box.text_frame
            p_hb = tf_hb.paragraphs[0]
            p_hb.text = "SMART INDIA HACKATHON 2026"
            p_hb.font.size = Pt(28)
            p_hb.font.bold = True
            p_hb.font.color.rgb = c_heading_blue
            p_hb.font.name = 'Georgia'
            p_hb.alignment = PP_ALIGN.CENTER

            # Sub-title
            subtitle_box = slide.shapes.add_textbox(Inches(2.0), Inches(1.4), Inches(8.0), Inches(0.6))
            tf_sb = subtitle_box.text_frame
            p_sb = tf_sb.paragraphs[0]
            p_sb.text = slide_data["title"]
            p_sb.font.size = Pt(22)
            p_sb.font.bold = True
            p_sb.font.color.rgb = c_black
            p_sb.font.name = 'Georgia'
            p_sb.alignment = PP_ALIGN.CENTER

            # Left lists parameters
            list_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.2))
            tf_list = list_box.text_frame
            tf_list.word_wrap = True
            
            for b_idx, bullet in enumerate(slide_data["bullets"]):
                p_b = tf_list.add_paragraph() if b_idx > 0 else tf_list.paragraphs[0]
                p_b.text = "•  " + bullet
                p_b.font.size = Pt(20)
                p_b.font.bold = True
                p_b.font.color.rgb = c_black
                p_b.font.name = 'Arial'
                p_b.space_after = Pt(14)

        elif slide_data["type"] == "solution":
            # Slide Header
            header_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.4), Inches(8.0), Inches(0.8))
            tf_hb = header_box.text_frame
            p_hb = tf_hb.paragraphs[0]
            p_hb.text = slide_data["title"]
            p_hb.font.size = Pt(26)
            p_hb.font.bold = True
            p_hb.font.color.rgb = c_black
            p_hb.font.name = 'Georgia'
            p_hb.alignment = PP_ALIGN.CENTER

            # Sub heading
            subheading_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.6))
            tf_sh = subheading_box.text_frame
            p_sh = tf_sh.paragraphs[0]
            p_sh.text = slide_data["subtitle"]
            p_sh.font.size = Pt(22)
            p_sh.font.bold = True
            p_sh.font.color.rgb = c_heading_blue
            p_sh.font.name = 'Georgia'

            # Bullet content
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.333), Inches(4.3))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            
            for b_idx, bullet in enumerate(slide_data["bullets"]):
                p_b = tf_body.add_paragraph() if b_idx > 0 else tf_body.paragraphs[0]
                p_b.text = "•  " + bullet
                p_b.font.size = Pt(18)
                p_b.font.color.rgb = c_black
                p_b.font.name = 'Arial'
                p_b.space_after = Pt(12)
                p_b.line_spacing = 1.1

        else:
            # Standard Slide title
            header_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.4), Inches(8.0), Inches(0.8))
            tf_hb = header_box.text_frame
            p_hb = tf_hb.paragraphs[0]
            p_hb.text = slide_data["title"]
            p_hb.font.size = Pt(26)
            p_hb.font.bold = True
            p_hb.font.color.rgb = c_black
            p_hb.font.name = 'Georgia'
            p_hb.alignment = PP_ALIGN.CENTER

            # Bullet content
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            
            for b_idx, bullet in enumerate(slide_data["bullets"]):
                p_b = tf_body.add_paragraph() if b_idx > 0 else tf_body.paragraphs[0]
                
                # Check for bullet indentation (two spaces representation)
                if bullet.startswith("  "):
                    p_b.text = "    " + bullet.strip()
                    p_b.font.size = Pt(16)
                    p_b.font.color.rgb = c_muted
                else:
                    p_b.text = "•  " + bullet.strip()
                    p_b.font.size = Pt(18)
                    p_b.font.color.rgb = c_black
                    
                p_b.font.name = 'Arial'
                p_b.space_after = Pt(8)
                p_b.line_spacing = 1.1

    # Save final file output
    prs.save("AICTE_M6_SIH_Submission.pptx")
    print("SIH Idea Submission PPT generated successfully as 'AICTE_M6_SIH_Submission.pptx'")

if __name__ == "__main__":
    create_sih_presentation()
