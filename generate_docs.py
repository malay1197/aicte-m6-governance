import sys
import os

try:
    from docx import Document
    from docx.shared import Inches as DocInches, Pt as DocPt
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
except ImportError:
    print("python-docx is not installed. Please install it using pip.")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Please install it using pip.")
    sys.exit(1)


def generate_docx():
    print("Generating Word Document...")
    doc = Document()
    
    # Title
    title = doc.add_heading(level=0)
    run = title.add_run("AICTE Security & Governance Platform\nModule M6: Audit, Attendance, Reports & Memory")
    run.font.name = 'Arial'
    run.bold = True
    
    # Section 1: Executive Summary
    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        "This document describes the design, architecture, and database layout of the Module M6 "
        "(Audit + Attendance + Reports + Memory) of the Secure AICTE Online Meeting & Governance Platform. "
        "The M6 module handles system audit trails, join/leave duration auditing, compliance reporting, "
        "institutional memory searching, and security notifications."
    )
    
    # Section 2: Core Responsibilities
    doc.add_heading("2. Core Module Responsibilities", level=1)
    p = doc.add_paragraph("The M6 module is responsible for the following governance operations:")
    p.paragraph_format.left_indent = DocInches(0.2)
    doc.add_paragraph("• Audit Trails: Detailed entry of logs including user logins, meetings, files, and permissions.", style='List Bullet')
    doc.add_paragraph("• Attendance Tracking: Calculating active join/leave durations and tag status (Present, Late, Absent, Left Early).", style='List Bullet')
    doc.add_paragraph("• Security Operations Center: Active threat monitoring, VPN exit node tracking, and IP quarantines.", style='List Bullet')
    doc.add_paragraph("• Parameterized Reports: Generation of summaries with cryptographic verification hashes.", style='List Bullet')
    doc.add_paragraph("• Institutional Memory Search: Full-Text Search (FTS) queries with role-based restrictions.", style='List Bullet')

    # Section 3: Database Schemas
    doc.add_heading("3. Database Schema Models", level=1)
    doc.add_paragraph(
        "The PostgreSQL DDL schema defines five primary tables to store M6 states. "
        "Below are the structural specifications for each table:"
    )
    
    # Table 1 specifications
    doc.add_heading("3.1 Table: audit_logs", level=2)
    table = doc.add_table(rows=6, cols=3)
    table.style = 'Light Shading Accent 1'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Column Name'
    hdr_cells[1].text = 'Data Type'
    hdr_cells[2].text = 'Description'
    
    rows = [
        ('id', 'UUID (PK)', 'Unique identifier of the log entry'),
        ('actor_username', 'VARCHAR(50)', 'Identifier of the user performing action'),
        ('action', 'VARCHAR(100)', 'Type of action (e.g. Permission Changed)'),
        ('severity_level', 'VARCHAR(15)', 'INFO, WARNING, CRITICAL levels'),
        ('ip_address', 'INET', 'Origin network address')
    ]
    for i, (name, dtype, desc) in enumerate(rows):
        row_cells = table.rows[i+1].cells
        row_cells[0].text = name
        row_cells[1].text = dtype
        row_cells[2].text = desc

    # Section 4: API Endpoint Specifications
    doc.add_heading("4. REST API Endpoint Specifications", level=1)
    doc.add_paragraph("The following REST interfaces are implemented in the Express backend routing layer:")
    
    endpoints = [
        ("GET /api/audit-logs", "Query compliance events with severity and user filters."),
        ("GET /api/attendance", "Fetch meeting participant attendance registers."),
        ("POST /api/reports/generate", "Compile reports and issue blockchain seals."),
        ("GET /api/memory/search", "FTS query search with role-based restriction layers."),
        ("PATCH /api/notifications/:id/read", "Mark system notification read.")
    ]
    for route, desc in endpoints:
        p = doc.add_paragraph()
        r = p.add_run(f"• {route}: ")
        r.bold = True
        p.add_run(desc)

    doc.save("AICTE_M6_Technical_Documentation.docx")
    print("Word Document generated successfully as 'AICTE_M6_Technical_Documentation.docx'")


def generate_pptx():
    print("Generating PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = PptInches(13.333) # 16:9 format
    prs.slide_height = PptInches(7.5)
    
    # Custom colors
    c_bg = RGBColor(10, 15, 29)       # Deep Navy #0A0F1D
    c_title = RGBColor(248, 250, 252) # White
    c_body = RGBColor(148, 163, 184)  # Muted slate
    c_indigo = RGBColor(96, 165, 250) # Light blue/indigo
    c_green = RGBColor(16, 185, 129)  # Emerald
    
    slides_data = [
        # Slide 1: Title
        {
            "is_title": True,
            "title": "Secure AICTE Online Meeting\n& Governance Platform",
            "subtitle": "Module M6: Audit + Attendance + Reports + Memory\nSmart India Hackathon Presentation Deck"
        },
        # Slide 2: Scope
        {
            "title": "Module M6: Scope & Responsibilities",
            "bullets": [
                "Audit logs: Security tracking of logins, permission changes, and document downloads.",
                "Meeting Attendance: Dynamic join/leave analysis and session connection durations.",
                "Admin Reports: Multi-parameter analytics generator with PDF/CSV export controls.",
                "SOC Threats Panel: Real-time warning mitigations and malicious IP quarantines.",
                "Institutional Search: Authorized search portal across meetings, decisions, and files."
            ]
        },
        # Slide 3: Architecture
        {
            "title": "M6 System Architecture & Integrations",
            "bullets": [
                "M1 Auth Integration: Ingests user identity contexts and enforces role clearances.",
                "M2 Meeting Integration: Receives participant connect/disconnect webhooks.",
                "M4 Blockchain Link: Anchors compliance hashes (SHA-256) of generated reports to blockchain ledger.",
                "M5 AI Transcription Link: Pulls summarized decisions and action-items to index memory.",
                "Central Compliance Engine: Stores audit logs, alerts, notifications, and indices."
            ]
        },
        # Slide 4: Database Design
        {
            "title": "Database Schema Design (PostgreSQL)",
            "bullets": [
                "audit_logs: Logs timestamps, actions, targets, client IPs, statuses, and severities.",
                "attendance: Tracks meeting participant logs, computing total minutes and status tags.",
                "compliance_reports: Stores report configurations and blockchain ledger hashes.",
                "notifications: Manages warning notifications, prioritization levels, and read indices.",
                "institutional_memory: Full-Text Search (FTS) catalog optimized with GIN indexing."
            ]
        },
        # Slide 5: REST API Specifications
        {
            "title": "REST API Endpoint Specifications",
            "bullets": [
                "GET /api/audit-logs : Queries compliance audit streams.",
                "GET /api/attendance : Fetches meeting participant registers.",
                "GET /api/reports : Loads reporting templates and historical hashes.",
                "POST /api/reports/generate : Compiles report data and calls M4 blockchain seals.",
                "GET /api/memory/search : Queries FTS index with role-based restriction checks.",
                "PATCH /api/notifications/:id/read : Marks alert notifications as read."
            ]
        },
        # Slide 6: Attendance Module Mechanics
        {
            "title": "Attendance Module Mechanics",
            "bullets": [
                "Integrates Jitsi webhooks to record user connecting and leaving events.",
                "Computes active connection time in minutes dynamically.",
                "Smart Attendance status classifications:",
                "  - Present: Logged in on-time and completed session.",
                "  - Late: Logged in 5 minutes past scheduled start.",
                "  - Left Early: Logged out before the meeting was formally adjourned.",
                "  - Absent: No connection logs detected in meeting duration."
            ]
        },
        # Slide 7: Security Operations Centre (SOC)",
        {
            "title": "Security Audits & SOC Remediation",
            "bullets": [
                "Logs compliance events across modules (Auth, Meetings, Files, Ledger).",
                "Color-coded warning notifications: INFO (Blue), WARNING (Orange), CRITICAL (Red).",
                "Critical alerts trigger warning badges on headers and admin panels.",
                "Remediation Modal: Allows compliance officers to quarantine IPs and sign off security cases."
            ]
        },
        # Slide 8: Compliance Reports & Blockchain Ledger
        {
            "title": "Reports Engine & Blockchain Sealing",
            "bullets": [
                "Parameterizes filters: meeting selection, dates, and report types.",
                "Compiles data into a simulated PDF preview layout directly inside the browser.",
                "Report payload hashed (SHA-256) and committed to M4 Blockchain ledger.",
                "Retrieves and embeds block hash signature as proof of compliance integrity."
            ]
        },
        # Slide 9: Authorized Institutional Memory
        {
            "title": "Authorized Institutional Memory",
            "bullets": [
                "Queries the FTS database. Example: searching 'budget decision' matches text.",
                "Relevance Ranking: Sorts and displays matching search scores (e.g. 98% Match).",
                "M1 Role Filter: Restricts search results. Users only see records matching their access role.",
                "Syncs transcripts (M5 AI) and decentralized block seals (M4 Ledger)."
            ]
        },
        # Slide 10: Conclusion & Deliverables
        {
            "title": "M6 Module Deliverables Summary",
            "bullets": [
                "Fully working React + Vite frontend prototype with dark-navy theme.",
                "Integrated Node.js + Express backend server serving REST endpoints.",
                "PostgreSQL schema scripts including FTS search vectors.",
                "Guided presenter deck allowing automated demonstration of compliance flows."
            ]
        }
    ]

    blank_layout = prs.slide_layouts[6]
    
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg
        
        if slide_data.get("is_title"):
            # Title slide text box
            title_box = slide.shapes.add_textbox(PptInches(1.0), PptInches(2.0), PptInches(11.333), PptInches(4.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            p1.text = slide_data["title"]
            p1.font.size = PptPt(44)
            p1.font.bold = True
            p1.font.color.rgb = c_title
            p1.font.name = 'Arial'
            
            p2 = tf.add_paragraph()
            p2.text = "\n" + slide_data["subtitle"]
            p2.font.size = PptPt(18)
            p2.font.color.rgb = c_indigo
            p2.font.name = 'Arial'
            
        else:
            # Header text box
            header_box = slide.shapes.add_textbox(PptInches(1.0), PptInches(0.6), PptInches(11.333), PptInches(1.2))
            tf_h = header_box.text_frame
            tf_h.word_wrap = True
            p_h = tf_h.paragraphs[0]
            p_h.text = slide_data["title"]
            p_h.font.size = PptPt(36)
            p_h.font.bold = True
            p_h.font.color.rgb = c_indigo
            p_h.font.name = 'Arial'
            
            # Content text box
            content_box = slide.shapes.add_textbox(PptInches(1.0), PptInches(2.0), PptInches(11.333), PptInches(4.8))
            tf_c = content_box.text_frame
            tf_c.word_wrap = True
            
            for i, bullet in enumerate(slide_data["bullets"]):
                p_c = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
                p_c.text = bullet
                p_c.font.size = PptPt(18)
                p_c.font.color.rgb = c_body
                p_c.font.name = 'Arial'
                p_c.space_after = PptPt(12)
                
            # Footer page numbering
            footer_box = slide.shapes.add_textbox(PptInches(11.0), PptInches(6.8), PptInches(1.5), PptInches(0.5))
            tf_f = footer_box.text_frame
            p_f = tf_f.paragraphs[0]
            p_f.text = f"Slide {idx + 1} / 10"
            p_f.font.size = PptPt(10)
            p_f.font.color.rgb = c_body
            p_f.font.name = 'Arial'

    prs.save("AICTE_M6_Presentation_10_Slides.pptx")
    print("PowerPoint Presentation generated successfully as 'AICTE_M6_Presentation_10_Slides.pptx'")


if __name__ == "__main__":
    generate_docx()
    generate_pptx()
